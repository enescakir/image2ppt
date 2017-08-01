#!/usr/local/bin/python
import sys, os
from pptx import Presentation
from pptx.util import Inches
import Tkinter
from Tkinter import *
import ScrolledText as tkst
import tkFileDialog as filedialog

DEFAULT_INPUT_PATH = os.path.expanduser("~/Desktop/")
DEFAULT_OUTPUT_PATH = os.path.expanduser("~/Desktop/")
DEFAULT_PPT_NAME = "image2ppt"
DEFAULT_SLIDE_WIDTH = 1920
DEFAULT_SLIDE_HEIGHT = 1080
DEFAULT_RESOLUTION = 72
SUPPORTED_FORMATS = ['.jpg', '.JPG', '.jpeg', '.JPEG', '.png', '.PNG' ]

ppt_name = DEFAULT_PPT_NAME
input_path = DEFAULT_INPUT_PATH
output_path = DEFAULT_OUTPUT_PATH
slide_width = DEFAULT_SLIDE_WIDTH
slide_height = DEFAULT_SLIDE_HEIGHT
resolution = DEFAULT_RESOLUTION

def pixelToEmu(pix):
    return int(int(pix) * 914400 / int(resolution))

def getInputFolder():
    global input_path
    input_path = filedialog.askdirectory()
    l5.config(text=input_path)

def getOutputFolder():
    global output_path
    output_path = filedialog.askdirectory()
    l6.config(text=output_path)

def log(text):
    tb.insert('insert', text + "\n")
    tb.see(Tkinter.END)

def generate():
    tb.delete('1.0', END)
    global ppt_name, resolution, slide_width, slide_height
    ppt_name = e1.get()
    resolution = int(e4.get())
    slide_width = pixelToEmu(e2.get())
    slide_height = pixelToEmu(e3.get())

    files = os.listdir(input_path)
    images = []
    for f in files:
        path = os.path.join(input_path, f)
        if os.path.isfile(path):
            filename, extension = os.path.splitext(f)
            if extension in SUPPORTED_FORMATS:
                images.append(path)

    log(str(len(images)) + " images found")

    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    for image in images:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(image, Inches(0) , Inches(0), slide_width,  slide_height)
        log("\"" + os.path.basename(image) + "\" is added")


    prs.save(os.path.join(output_path, ppt_name + '.pptx'))
    log(ppt_name + '.pptx is saved')

root = Tk()
root.resizable(width=False, height=False)
root.minsize(width=500, height=420)
root.wm_title("image2ppt")
root.attributes('-topmost', 1)
root.after(1000, lambda: root.attributes('-topmost', 0))

l0 = Label(root, text="image2ppt", font = "Helvetica 24 bold")
l1 = Label(root, text="Presentation's name:", anchor=W)
l2 = Label(root, text="Image width in pixels:", anchor=W)
l3 = Label(root, text="Image height in pixels:", anchor=W)
l4 = Label(root, text="Image resolution in ppi:", anchor=W)

l0.place(x = 150, y = 10, width=200, height=40)
l1.place(x = 100, y = 60, width=160, height=30)
l2.place(x = 100, y = 100, width=160, height=30)
l3.place(x = 100, y = 140, width=160, height=30)
l4.place(x = 100, y = 180, width=160, height=30)

e1 = Entry(root)
e2 = Entry(root)
e3 = Entry(root)
e4 = Entry(root)

e1.insert(0, DEFAULT_PPT_NAME)
e2.insert(0, DEFAULT_SLIDE_WIDTH)
e3.insert(0, DEFAULT_SLIDE_HEIGHT)
e4.insert(0, DEFAULT_RESOLUTION)

e1.place(x = 280, y = 60, width=120, height=30)
e2.place(x = 280, y = 100, width=120, height=30)
e3.place(x = 280, y = 140, width=120, height=30)
e4.place(x = 280, y = 180, width=120, height=30)

l5 = Label(root, text=DEFAULT_INPUT_PATH, anchor=W)
b1 = Button(root, text='Choose Images', command=getInputFolder)
l5.place(x = 20, y = 220, width=330, height=30)
b1.place(x = 360, y = 220, width=120, height=30)

l6 = Label(root, text=DEFAULT_OUTPUT_PATH, anchor=W)
b2 = Button(root, text='Choose Output', command=getOutputFolder)
l6.place(x = 20, y = 260, width=330, height=30)
b2.place(x = 360, y = 260, width=120, height=30)

b3 = Button(root, text='Generate Presentation', command=generate)
b3.place(x = 100, y = 300, width=300)

tb = tkst.ScrolledText(
    master = root,
    wrap   = 'word',  # wrap text at full words only
    width  = 25,      # characters
    height = 10,      # text lines
    # state=DISABLED,
    borderwidth=2,
    relief="groove"
)
tb.place(x = 50, y = 340, width=400, height=60)

root.mainloop()
